import streamlit as st
import google.generativeai as genai
import PyPDF2
import docx
from PIL import Image
import pandas as pd
import io
import os
import re
import smtplib
import ssl
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
from docx.enum.text import WD_COLOR_INDEX

# ==========================================
# CẤU HÌNH BACKEND: Khai báo API Key ở đây
# ==========================================
# Bạn hãy thay thế chuỗi bên dưới bằng API Key thật của bạn.
# Tuyệt đối không để lộ mã này lên GitHub công khai.
GOOGLE_API_KEY = "PASTE_YOUR_API_KEY_HERE"

# Ưu tiên cấu hình từ Streamlit Secrets, nếu không có thì lấy trực tiếp từ biến trên
try:
    API_KEY = st.secrets.get("GEMINI_API_KEY", GOOGLE_API_KEY)
except Exception:
    API_KEY = GOOGLE_API_KEY

def configure_genai():
    if API_KEY and API_KEY != "PASTE_YOUR_API_KEY_HERE":
        genai.configure(api_key=API_KEY)
        return True
    return False

# ==========================================
# THIẾT LẬP GIAO DIỆN (UI)
# ==========================================
st.set_page_config(page_title="Công cụ số hóa", page_icon="📝", layout="wide")

# Áp dụng Custom CSS cho tông màu Xanh đậm - Trắng và Footer
st.markdown("""
    <style>
    /* Chỉnh màu chữ tiêu đề chính */
    .stApp {
        background-color: #FFFFFF;
    }
    h1, h2, h3 {
        color: #003366 !important; /* Xanh dương đậm */
    }
    /* Tùy chỉnh Sidebar */
    [data-testid="stSidebar"] {
        background-color: #F0F4F8;
    }
    /* Chỉnh sửa layout Markdown Table cho đẹp */
    table {
        width: 100%;
        border-collapse: collapse;
    }
    th {
        background-color: #004080;
        color: white;
        text-align: left;
        padding: 8px;
    }
    td {
        border: 1px solid #ddd;
        padding: 8px;
    }
    tr:nth-child(even) {
        background-color: #f2f2f2;
    }
    /* Footer */
    .footer {
        position: fixed;
        left: 0;
        bottom: 0;
        width: 100%;
        background-color: transparent;
        color: gray;
        text-align: center;
        padding: 10px;
        font-size: 14px;
        z-index: 100;
    }
    </style>
""", unsafe_allow_html=True)

# ==========================================
# CÁC HÀM XỬ LÝ ĐỌC FILE
# ==========================================
def extract_text_from_pdf(file):
    text = ""
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    except Exception as e:
        raise Exception(f"Không thể đọc file PDF (có thể là file scan hoặc bị lỗi): {e}")
    return text

def extract_text_from_docx(file):
    try:
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise Exception(f"Không thể đọc file Word: {e}")

def extract_text_from_pptx(file):
    from pptx import Presentation
    text = ""
    try:
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        raise Exception(f"Không thể đọc file PowerPoint: {e}")
    return text

def convert_doc_to_docx_win32(uploaded_file):
    import tempfile
    import os
    import win32com.client as win32
    import pythoncom
    
    # Save uploaded file bytes to a temporary .doc file
    temp_dir = tempfile.gettempdir()
    input_path = os.path.join(temp_dir, "temp_uploaded_giaoan.doc")
    output_path = os.path.splitext(input_path)[0] + ".docx"
    
    # Remove existing files if any
    if os.path.exists(input_path):
        try: os.remove(input_path)
        except: pass
    if os.path.exists(output_path):
        try: os.remove(output_path)
        except: pass
        
    uploaded_file.seek(0)
    with open(input_path, "wb") as f:
        f.write(uploaded_file.read())
        
    try:
        pythoncom.CoInitialize()
        word = win32.gencache.EnsureDispatch('Word.Application')
        word.Visible = False
        doc = word.Documents.Open(input_path)
        doc.SaveAs(output_path, FileFormat=12) # 12 = wdFormatXMLDocument
        doc.Close()
        word.Quit()
        
        # Read the newly created docx file
        with open(output_path, "rb") as docx_file:
            text = extract_text_from_docx(docx_file)
            
        return text
    except Exception as e:
        raise Exception(f"Lỗi khi chuyển đổi và đọc file .doc: {e}. Đảm bảo máy tính chạy ứng dụng đã cài đặt MS Word.")
    finally:
        if os.path.exists(input_path):
            try: os.remove(input_path)
            except: pass
        if os.path.exists(output_path):
            try: os.remove(output_path)
            except: pass

def convert_ppt_to_pptx_win32(uploaded_file):
    import tempfile
    import os
    import win32com.client as win32
    import pythoncom
    
    temp_dir = tempfile.gettempdir()
    input_path = os.path.join(temp_dir, "temp_uploaded_giaoan.ppt")
    output_path = os.path.splitext(input_path)[0] + ".pptx"
    
    if os.path.exists(input_path):
        try: os.remove(input_path)
        except: pass
    if os.path.exists(output_path):
        try: os.remove(output_path)
        except: pass
        
    uploaded_file.seek(0)
    with open(input_path, "wb") as f:
        f.write(uploaded_file.read())
        
    try:
        pythoncom.CoInitialize()
        powerpoint = win32.Dispatch("Powerpoint.Application")
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        presentation.SaveAs(output_path, 24) # 24 = PpSaveAsFileType.ppSaveAsOpenXMLPresentation
        presentation.Close()
        powerpoint.Quit()
        
        with open(output_path, "rb") as pptx_file:
            text = extract_text_from_pptx(pptx_file)
            
        return text
    except Exception as e:
        raise Exception(f"Lỗi khi chuyển đổi và đọc file .ppt: {e}. Đảm bảo máy tính chạy ứng dụng đã cài đặt MS PowerPoint.")
    finally:
        if os.path.exists(input_path):
            try: os.remove(input_path)
            except: pass
        if os.path.exists(output_path):
            try: os.remove(output_path)
            except: pass

# Hàm chuyển đổi Markdown Table sang DataFrame của Pandas
def markdown_table_to_df(markdown_str):
    # Tìm tất cả các dòng chứa ký tự '|' báo hiệu bảng
    lines = markdown_str.strip().split('\n')
    table_lines = [line for line in lines if '|' in line]
    
    if not table_lines:
        return None
        
    # Xử lý tiêu đề (dòng đầu tiên)
    header_line = table_lines[0]
    headers = [col.strip() for col in header_line.split('|') if col.strip()]
    
    # Xử lý các dòng dữ liệu (bỏ qua dòng phân cách ---|--- thường là dòng số 2)
    data = []
    for line in table_lines[1:]:
        # Bỏ qua dòng format ----
        if set(line.replace('|', '').replace('-', '').replace(' ', '').replace(':', '')) == set():
            continue
        cols = [col.strip() for col in line.split('|')[1:-1]] # Bỏ cột rỗng ở đầu và cuối do split
        if len(cols) > 0:
            # Cắt hoặc padding thêm nếu số cột không khớp
            if len(cols) > len(headers):
                cols = cols[:len(headers)]
            elif len(cols) < len(headers):
                cols = cols + [""] * (len(headers) - len(cols))
            
            # Cấu hình xuống dòng thực thụ cho Excel
            cols = [col.replace("<br>", "\n").replace("<br/>", "\n").replace("<br />", "\n") for col in cols]
            data.append(cols)
            
    if headers and data:
        return pd.DataFrame(data, columns=headers)
    return None

# ==========================================
# GIAO DIỆN CHÍNH
# ==========================================
st.title("🏛️ Công cụ số hóa")
st.markdown("**phục vụ công việc chuyển đổi số**")
st.divider()

tab_noi_tru, tab_giao_an = st.tabs(["📋 Quản lý nội trú", "📝 Soạn Giáo Án Năng Lực Số"])

# ------------------------------------------
# TAB 1: QUẢN LÝ NỘI TRÚ (Xử lý văn bản cũ)
# ------------------------------------------
with tab_noi_tru:
    col_left, col_right = st.columns([1, 2])
    
    with col_left:
        st.header("📂 Tải Văn Bản")
        uploaded_file = st.file_uploader(
            "Kéo thả hoặc dán file vào đây", 
            type=["pdf", "docx", "png", "jpg", "jpeg"]
        )
        st.markdown("---")
        st.markdown("""
        **✅ Hướng dẫn sử dụng:**
        1. Tải lên công văn, kế hoạch (File Word, PDF) hoặc ảnh chụp công văn có dấu đỏ.
        2. Đợi hệ thống AI đọc và xử lý.
        3. Nhận bảng công việc đã được bóc tách tự động.
        4. Tải file Excel về máy để lưu minh chứng theo dõi.
        """)
        
    with col_right:

        PROMPT_TEXT = """Đóng vai một Hiệu trưởng / Quản lý hành chính trường học. Hãy đọc văn bản chỉ đạo sau và bóc tách thông tin thành một bảng nghiêm ngặt. 
Bảng phải gồm chính xác 4 cột:
1. Tóm tắt Nội dung chính (Ngắn gọn 2-3 câu).
2. Đối tượng thực hiện (Ghi đích danh: GV Ngữ văn, Lịch sử, Ban giám hiệu, Bảo vệ...).
3. Hành động cần làm (Liệt kê gạch đầu dòng các công việc cụ thể. BẮT BUỘC dùng thẻ HTML <br> để xuống dòng giữa các gạch đầu dòng để giao diện hiển thị đẹp mắt).
4. Hạn hoàn thành (Rút trích ngày tháng, nếu văn bản không ghi thì điền 'Theo tiến độ chung').
Trả về kết quả 100% dưới dạng Markdown Table để tôi hiển thị lên web.
"""

        if uploaded_file is not None:
            if not configure_genai():
                st.error("⚠️ LỖI: Chưa cấu hình GOOGLE_API_KEY ở backend. Vui lòng kiểm tra mã nguồn (app.py) hoặc cấu hình Streamlit Secrets.")
            else:
                st.info(f"Đang phân tích tài liệu: **{uploaded_file.name}**...")
                
                try:
                    # Lấy danh sách model khả dụng
                    available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
                    if not available_models:
                        raise Exception("API Key của bạn không có quyền truy cập vào bất kỳ mô hình Gemini nào hỗ trợ tạo nội dung.")
                        
                    # Ưu tiên chọn gemini-1.5-flash hoặc gemini-2.5-flash nếu có, nếu không thì lấy model đầu tiên
                    selected_model = available_models[0]
                    for m_name in available_models:
                        if "1.5-flash" in m_name or "2.5-flash" in m_name:
                            selected_model = m_name
                            break
                            
                    model = genai.GenerativeModel(selected_model)
                    
                    with st.spinner(f"AI ({selected_model}) đang bóc tách dữ liệu... Vui lòng đợi trong giây lát."):
                        response = None
                        
                        file_ext = uploaded_file.name.split('.')[-1].lower()
                        
                        # Xử lý ảnh (Gửi thẳng file ảnh qua Vision model)
                        if file_ext in ['png', 'jpg', 'jpeg']:
                            image = Image.open(uploaded_file)
                            st.image(image, caption="Ảnh chụp công văn tải lên", width=300)
                            response = model.generate_content([PROMPT_TEXT, image])
                            
                        # Xử lý text từ PDF hoặc DOCX
                        else:
                            text_content = ""
                            if file_ext == "pdf":
                                text_content = extract_text_from_pdf(uploaded_file)
                            elif file_ext == "docx":
                                text_content = extract_text_from_docx(uploaded_file)
                            
                            if not text_content.strip():
                                st.warning("⚠️ Không tìm thấy chữ trong văn bản. Nếu đây là PDF dạng scan (văn bản chụp hình), vui lòng chuyển sang file ảnh (.png, .jpg) để upload lại.")
                            else:
                                full_prompt = PROMPT_TEXT + "\\n\\nNội dung văn bản:\\n" + text_content
                                response = model.generate_content(full_prompt)
                        
                        # Render kết quả
                        if response:
                            st.success("✅ Đã bóc tách thành công!")
                            
                            st.subheader("📊 Bảng Phân công Công việc")
                            markdown_result = response.text
                            
                            # Hiện bảng lên màn hình và cho phép render thẻ HTML <br>
                            st.markdown(markdown_result, unsafe_allow_html=True)
                            
                            # Xử lý xuất Excel
                            df = markdown_table_to_df(markdown_result)
                            if df is not None:
                                # Ghi Dataframe ra bộ nhớ đệm (buffer) để tạo file Excel tải xuống
                                output = io.BytesIO()
                                with pd.ExcelWriter(output, engine='openpyxl') as writer:
                                    df.to_excel(writer, index=False, sheet_name='Phan_Cong')
                                
                                excel_data = output.getvalue()
                                
                                st.markdown("---")
                                col1, col2, col3 = st.columns([1, 2, 1])
                                with col2:
                                    st.download_button(
                                        label="📥 Tải xuống Bảng Phân công (Excel)",
                                        data=excel_data,
                                        file_name=f"Ban_Phan_Cong_{uploaded_file.name}.xlsx",
                                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                        use_container_width=True
                                    )
                            else:
                                st.warning("⚠️ AI trả về kết quả nhưng không nằm trong định dạng bảng chuẩn nên không thể tạo file Excel. Xin thử lại với tư duy khác của AI.")
                                
                except Exception as e:
                    st.error(f"❌ Xảy ra lỗi trong quá trình xử lý: {str(e)}")
        else:
            st.info("👈 Vui lòng tải tài liệu lên từ bảng bên trái để bắt đầu bóc tách phân công công việc.")

# ------------------------------------------
# TAB 2: SOẠN GIÁO ÁN NĂNG LỰC SỐ
# ------------------------------------------
with tab_giao_an:
    st.header("📝 Soạn Giáo Án Năng Lực Số")
    
    col_em1, col_em2 = st.columns([1, 1])
    with col_em1:
        email_nhan = st.text_input("Nhập Email giáo viên nhận giáo án:", placeholder="username@gmail.com")
    with col_em2:
        gmail_app_password = st.text_input("Nhập Gmail App Password của giáo viên:", type="password", placeholder="Mật khẩu ứng dụng 16 ký tự...")
        
    st.markdown("---")
    
    # Sử dụng Session State để lưu văn bản trích xuất từ file giáo án tải lên
    if "giao_an_text" not in st.session_state:
        st.session_state.giao_an_text = ""
        
    uploaded_giao_an_file = st.file_uploader(
        "Tải lên giáo án cũ của bạn (Hỗ trợ PDF, Word .doc/.docx, PowerPoint .ppt/.pptx):",
        type=["pdf", "docx", "pptx", "doc", "ppt"]
    )
    
    if uploaded_giao_an_file is not None:
        file_ext = uploaded_giao_an_file.name.split('.')[-1].lower()
        try:
            with st.spinner("Đang trích xuất văn bản từ tệp tin tải lên..."):
                if file_ext == "pdf":
                    extracted_text = extract_text_from_pdf(uploaded_giao_an_file)
                elif file_ext == "docx":
                    extracted_text = extract_text_from_docx(uploaded_giao_an_file)
                elif file_ext == "pptx":
                    extracted_text = extract_text_from_pptx(uploaded_giao_an_file)
                elif file_ext == "doc":
                    extracted_text = convert_doc_to_docx_win32(uploaded_giao_an_file)
                elif file_ext == "ppt":
                    extracted_text = convert_ppt_to_pptx_win32(uploaded_giao_an_file)
                else:
                    extracted_text = ""
            
            if extracted_text.strip():
                st.session_state.giao_an_text = extracted_text
                st.success("✅ Trích xuất văn bản giáo án thành công! Nội dung đã được điền vào ô bên dưới.")
            else:
                st.warning("⚠️ Không tìm thấy hoặc không thể đọc được nội dung văn bản từ tệp tin này.")
        except Exception as read_err:
            st.error(f"❌ Lỗi khi đọc file: {str(read_err)}")
                
    giao_an_cu = st.text_area(
        "Nội dung giáo án Ngữ văn cũ:",
        value=st.session_state.giao_an_text,
        height=300,
        placeholder="Dán giáo án cũ tại đây hoặc chọn file tải lên ở trên để hệ thống tự động đọc nội dung..."
    )
    
    submit_btn = st.button("TÍCH HỢP KHUNG NĂNG LỰC SỐ VÀ CẤP PHÁT GIÁO ÁN", use_container_width=True)
    
    if submit_btn:
        if not email_nhan.strip():
            st.warning("⚠️ Vui lòng nhập Email giáo viên nhận giáo án!")
        elif not gmail_app_password.strip():
            st.warning("⚠️ Vui lòng nhập Gmail App Password của giáo viên!")
        elif not giao_an_cu.strip():
            st.warning("⚠️ Vui lòng nhập hoặc dán nội dung giáo án cũ!")
        else:
            try:
                # Cấu hình Gemini bằng API Key hệ thống mặc định
                if not configure_genai():
                    st.error("⚠️ LỖI: Chưa cấu hình GOOGLE_API_KEY ở backend. Vui lòng kiểm tra mã nguồn (app.py) hoặc cấu hình Streamlit Secrets.")
                else:
                    system_instruction = """Bạn là chuyên gia thẩm định và xây dựng chương trình giáo dục phổ thông môn Ngữ văn. Hãy đọc giáo án cũ được cung cấp và tiến hành nâng cấp, tích hợp Khung năng lực số cho người học theo Thông tư số 02/2025/TT-BGDĐT và Khung giáo dục AI theo Quyết định số 3439/QĐ-BGDĐT của Bộ Giáo dục và Đào tạo Việt Nam.
Yêu cầu tích hợp:
+ Giữ vững cấu trúc kiến thức đặc trưng của thể loại văn học (Sử thi, Thần thoại, Thơ trữ tình, Bi kịch, Tiểu thuyết...) theo đúng phân phối chương trình Sách giáo khoa Kết nối tri thức.
+ Đối với cấp THPT, thiết kế các hoạt động học tập tương tác số đạt Mức độ thành thạo Bậc 5 (Nâng cao 1). Học sinh phải đóng vai trò tự chủ: tự khai thác dữ liệu, sử dụng AI tạo sinh có trách nhiệm để kiểm chứng thông tin, tự tạo lập nội dung số đa phương thức hoặc thực hiện trách nhiệm công dân trong môi trường số.
+ Lồng ghép các mục tiêu số này vào mục 'Mục tiêu bài học (Về năng lực số)' và triển khai các hoạt động cụ thể của học sinh trong mục 'Tiến trình dạy học'.
+ Bắt buộc đánh dấu bắt đầu và kết thúc của mỗi đoạn văn bản được tích hợp mới bằng cụm từ chính xác: [DIGITAL_START] và [DIGITAL_END] để hệ thống xử lý hậu kỳ bôi vàng."""
                    
                    with st.spinner("AI đang phân tích và tích hợp Khung năng lực số vào giáo án..."):
                        model = genai.GenerativeModel("gemini-1.5-flash", system_instruction=system_instruction)
                        response = model.generate_content(f"Nội dung giáo án cũ:\n\n{giao_an_cu}")
                    
                    if response and response.text:
                        st.success("✅ Đã tích hợp Khung năng lực số thành công!")
                        
                        # Hiển thị kết quả lên màn hình
                        st.subheader("📊 Giáo án đã được nâng cấp tích hợp")
                        st.markdown(response.text)
                        
                        # Logic đóng gói file Word
                        with st.spinner("Đang đóng gói giáo án vào file Word..."):
                            doc = docx.Document()
                            paragraphs = response.text.split('\n')
                            for p_text in paragraphs:
                                is_highlighted = "[DIGITAL_START]" in p_text and "[DIGITAL_END]" in p_text
                                clean_text = p_text.replace("[DIGITAL_START]", "").replace("[DIGITAL_END]", "")
                                
                                p = doc.add_paragraph()
                                run = p.add_run(clean_text)
                                if is_highlighted:
                                    run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                            
                            doc.save("Giao_An_Nang_Luc_So.docx")
                        
                        # Logic gửi Email SMTP Gmail
                        msg = MIMEMultipart()
                        msg["From"] = email_nhan
                        msg["To"] = email_nhan
                        msg["Subject"] = "[Smart App] Giáo Án Ngữ Văn Tích Hợp Khung Năng Lực Số Hoàn Thiện"
                        
                        body = "Kính gửi Thầy/Cô,\n\nĐây là Giáo án Ngữ văn tích hợp Khung năng lực số đã được tự động thẩm định và nâng cấp hoàn thiện bởi AI.\n\nCác phần tích hợp mới đã được bôi vàng trong tài liệu đính kèm.\n\nTrân trọng,\nSmart App"
                        msg.attach(MIMEText(body, "plain", "utf-8"))
                        
                        filename = "Giao_An_Nang_Luc_So.docx"
                        with open(filename, "rb") as attachment:
                            part = MIMEBase("application", "octet-stream")
                            part.set_payload(attachment.read())
                            encoders.encode_base64(part)
                            part.add_header(
                                "Content-Disposition",
                                f"attachment; filename={filename}",
                            )
                            msg.attach(part)
                        
                        with st.spinner("Đang gửi email đính kèm giáo án qua máy chủ SMTP Gmail..."):
                            context = ssl.create_default_context()
                            with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=context) as server:
                                server.login(email_nhan, gmail_app_password)
                                server.sendmail(email_nhan, email_nhan, msg.as_string())
                            st.success(f"🚀 Đã gửi email đính kèm giáo án đến hộp thư **{email_nhan}** thành công!")
            except Exception as e:
                st.error(f"❌ Xảy ra lỗi trong quá trình xử lý hoặc gửi email: {str(e)}")

# ==========================================
# FOOTER
# ==========================================
st.markdown('<div class="footer">© 2026 Bản quyền thuộc về Đỗ Viết Cường - Trường PTDTNT Cao Lộc</div>', unsafe_allow_html=True)
